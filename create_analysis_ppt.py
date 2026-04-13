import pandas as pd
import numpy as np
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

file_path = "/root/.openclaw/media/inbound/PH新版日报2026年3月1-19日---54ca058d-ce20-4a34-8dba-f299499485d5.xlsx"
sheets = ["WP_ro", "GJP_ro", "SBET_ro", "Sugarplay_ro"]
output_ppt = "/root/.openclaw-tcs/workspace/PH_Platform_Analysis_Report.pptx"

def create_ppt():
    all_data = []
    for sheet in sheets:
        df = pd.read_excel(file_path, sheet_name=sheet, header=1)
        df.columns = [str(c).replace('\n', '').strip() for c in df.columns]
        mapping = {'日期': 'Date', '总登录人数': 'Logins', '总新注册人数': 'Regs', '总首充人数': 'FirstDeps', '总存款人数': 'DepUsers', '总存款金额': 'Amount', '总存款额': 'Amount'}
        rename_map = {k: v for k, v in mapping.items() if k in df.columns}
        if 'Amount' not in rename_map.values():
             df['Amount'] = df.iloc[:, 5]
             rename_map['Amount'] = 'Amount'
        df = df.rename(columns=rename_map)
        df['Date'] = pd.to_datetime(df['Date'], errors='coerce')
        df = df[df['Date'].notnull()]
        df = df[(df['Date'].dt.month == 3) & (df['Date'].dt.day <= 19)]
        for col in ['Logins', 'Regs', 'FirstDeps', 'DepUsers', 'Amount']:
            df[col] = pd.to_numeric(df[col], errors='coerce').fillna(0)
        df = df[df['Logins'].apply(lambda x: x == int(x))].drop_duplicates('Date')
        df['Platform'] = sheet.replace('_ro', '')
        all_data.append(df[['Date', 'Platform', 'Logins', 'Regs', 'FirstDeps', 'DepUsers', 'Amount']])

    df_full = pd.concat(all_data)
    
    # PPT Generation
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "PH 平台 3月运营深度分析报告"
    subtitle.text = "数据截止日期: 2026-03-19\n生成者: Tone (AI 助手)"

    # Market Share Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "市场规模与份额分析"
    body = slide.shapes.placeholders[1]
    share = df_full.groupby('Platform')['Amount'].sum() / df_full['Amount'].sum() * 100
    total_amt = df_full['Amount'].sum()
    text = f"全平台总存款金额: {total_amt:,.2f}\n\n"
    for p, s in share.items():
        text += f"• {p}: 份额 {s:.1f}%\n"
    body.text = text

    # Efficiency Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "运营效率与用户价值"
    body = slide.shapes.placeholders[1]
    efficiency = df_full.groupby('Platform').agg({'Regs': 'sum', 'FirstDeps': 'sum', 'DepUsers': 'sum', 'Amount': 'sum'})
    efficiency['ConvRate'] = (efficiency['FirstDeps'] / efficiency['Regs'] * 100).round(1)
    efficiency['ARPU'] = (efficiency['Amount'] / efficiency['DepUsers']).round(2)
    text = "核心效率指标:\n"
    for p, r in efficiency.iterrows():
        text += f"• {p}: 转化率 {r['ConvRate']}% | ARPU {r['ARPU']}\n"
    body.text = text

    # Growth Slide (Latest)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "最新增长动态 (3月19日)"
    body = slide.shapes.placeholders[1]
    latest_date = df_full['Date'].max()
    prev_date = latest_date - pd.Timedelta(days=1)
    l_data = df_full[df_full['Date'] == latest_date].set_index('Platform')
    p_data = df_full[df_full['Date'] == prev_date].set_index('Platform')
    growth = l_data['Amount'] - p_data['Amount']
    text = "对比昨日存款金额增长:\n"
    for p, g in growth.dropna().items():
        text += f"• {p}: {g:+.2f}\n"
    body.text = text

    # Strategy Slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "立体化调整建议"
    body = slide.shapes.placeholders[1]
    text = "1. WP: 流量大户，应加强转化链路，拉升ARPU值。\n"
    text += "2. GJP: 维持规模优势，优化新用户留存策略。\n"
    text += "3. SBET: 模式极佳，尝试寻找更高质量的流量入口。\n"
    text += "4. Sugarplay: 效率漏斗损耗严重，需紧急优化注册路径。"
    body.text = text

    prs.save(output_ppt)
    print(f"SUCCESS: {output_ppt}")

if __name__ == "__main__":
    create_ppt()
